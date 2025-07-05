import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from zipfile import BadZipFile
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pathlib import Path
import uuid
import io
import re
import tempfile
import os

# Load the API key securely from environment variable
api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=api_key)

# Default theme configurations
DEFAULT_THEMES = {
    "Professional Blue": {
        "background": RGBColor(12, 35, 64),
        "title_color": RGBColor(255, 255, 255),
        "text_color": RGBColor(200, 200, 200),
        "accent": RGBColor(0, 112, 192),
        "title_font": "Calibri",
        "text_font": "Calibri"
    },
    "Modern Green": {
        "background": RGBColor(22, 82, 66),
        "title_color": RGBColor(255, 255, 255),
        "text_color": RGBColor(220, 220, 220),
        "accent": RGBColor(76, 175, 80),
        "title_font": "Arial",
        "text_font": "Arial"
    },
    "Light Corporate": {
        "background": RGBColor(255, 255, 255),
        "title_color": RGBColor(13, 71, 161),
        "text_color": RGBColor(33, 33, 33),
        "accent": RGBColor(25, 118, 210),
        "title_font": "Segoe UI",
        "text_font": "Segoe UI"
    },
    "Dark Tech": {
        "background": RGBColor(33, 33, 33),
        "title_color": RGBColor(0, 200, 255),
        "text_color": RGBColor(200, 200, 200),
        "accent": RGBColor(0, 150, 255),
        "title_font": "Consolas",
        "text_font": "Consolas"
    }
}

def hex_to_rgb(hex_color):
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(*tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4)))

def extract_theme_from_pptx(uploaded_file):
    """Extract theme colors and fonts from an uploaded PowerPoint file"""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(uploaded_file.getbuffer())
        tmp_file_path = tmp_file.name
    
    prs = Presentation(tmp_file_path)
    theme = {
        "background": RGBColor(255, 255, 255),  # Default white if not found
        "title_color": RGBColor(0, 0, 0),       # Default black
        "text_color": RGBColor(0, 0, 0),        # Default black
        "accent": RGBColor(79, 129, 189),       # Default blue
        "title_font": "Calibri",
        "text_font": "Calibri",

        # Header and footer
        "header": None,
        "footer": None,
        "slide_number": False,
        
        # Title slide specific
        "title_slide_layout": None,
        "title_slide_background": None,
        "title_slide_title_style": None,
        "title_slide_subtitle_style": None,
        
        "template_path": tmp_file_path  # Store the template path for later use
    }
    
    try:
        # Get colors from master slide
        slide_master = prs.slide_master
        
        # Handle background color
        if hasattr(slide_master.background, 'fill'):
            fill = slide_master.background.fill
            if fill.type == 1 and fill.fore_color and fill.fore_color.rgb:  # Solid fill
                theme["background"] = fill.fore_color.rgb
        
        # Try to get title and text colors from placeholders
        for shape in slide_master.shapes:
            if shape.has_text_frame and shape.text.strip():
                try:
                    font = shape.text_frame.paragraphs[0].font
                    color = font.color.rgb if font.color and font.color.rgb else RGBColor(0, 0, 0)
                    name = font.name or "Calibri"

                    if "title" in shape.name.lower():
                        theme["title_color"] = color
                        theme["title_font"] = name
                    else:
                        theme["text_color"] = color
                        theme["text_font"] = name
                except:
                    continue

        # 2. Extract header and footer information
        for layout in slide_master.slide_layouts:
            for shape in layout.shapes:
                if not shape.has_text_frame:
                    continue
                try:
                    placeholder_format = getattr(shape, "placeholder_format", None)
                    if placeholder_format:
                        fmt_type = placeholder_format.type
                        font = shape.text_frame.paragraphs[0].font
                        color = font.color.rgb if font.color and font.color.rgb else RGBColor(0, 0, 0)
                        name = font.name or "Calibri"
                        size = font.size
        
                        details = {
                            "text": shape.text_frame.text,
                            "font": name,
                            "size": size,
                            "color": color,
                            "position": (shape.left, shape.top),
                            "size_dim": (shape.width, shape.height)
                        }
        
                        if fmt_type == 12:
                            theme["footer"] = details
                        elif fmt_type == 11:
                            theme["header"] = details
                        elif fmt_type == 13:
                            theme["slide_number"] = True
                except:
                    continue

        
        # 3. Extract title slide specific styles
        title_layout = None
        for layout in slide_master.slide_layouts:
            if "title" in layout.name.lower():
                title_layout = layout
                break
        
        if title_layout:
            theme["title_slide_layout"] = title_layout
            # Get background
            if hasattr(title_layout.background, 'fill'):
                theme["title_slide_background"] = title_layout.background.fill.fore_color.rgb
            
            # Get title and subtitle styles
            for shape in title_layout.shapes:
                if shape.has_text_frame:
                    try:
                        if "title" in shape.name.lower():
                            theme["title_slide_title_style"] = {
                                "font": shape.text_frame.paragraphs[0].font.name,
                                "size": shape.text_frame.paragraphs[0].font.size,
                                "color": shape.text_frame.paragraphs[0].font.color.rgb,
                                "position": (shape.left, shape.top),
                                "alignment": shape.text_frame.paragraphs[0].alignment
                            }
                        elif "subtitle" in shape.name.lower():
                            theme["title_slide_subtitle_style"] = {
                                "font": shape.text_frame.paragraphs[0].font.name,
                                "size": shape.text_frame.paragraphs[0].font.size,
                                "color": shape.text_frame.paragraphs[0].font.color.rgb,
                                "position": (shape.left, shape.top),
                                "alignment": shape.text_frame.paragraphs[0].alignment
                            }
                    except:
                        continue
    
        
        # Try to get accent color from first shape with fill
        for shape in slide_master.shapes:
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type == 1:  # Solid fill
                        theme["accent"] = shape.fill.fore_color.rgb
                        break
                except:
                    continue
    
    except Exception as e:
        st.warning(f"Couldn't fully extract theme: {str(e)}. Using default colors where needed.")
    
    return theme

def generate_slide_content(topic, slide_count):
    model = genai.GenerativeModel('gemini-2.0-flash')
    prompt = f"""Create a comprehensive presentation on '{topic}' with exactly {slide_count} slides. 
    For each slide, provide:
    1. A clear title in [Title:] format
    2. 3-5 detailed bullet points in [Content:] format (each point should be 2-3 lines/40-60 words)
    3. Optional speaker notes in [Notes:] format
    4. Layout suggestion in [Layout:] format (title-only, title-content, two-column, section-header)
    
    Structure your response like this:
    [Title:] Slide Title
    [Layout:] title-content
    [Content:]
    - Main Point 1: Detailed explanation spanning 1-2 lines with supporting information that provides context and value to the audience. This makes each point substantial.
    - Main Point 2: Another complete thought with sufficient detail to stand alone as a mini-paragraph, giving the audience concrete information they can use.
    - Main Point 3: Final point with enough depth to be meaningful, typically consisting of 2-3 sentences that develop a complete idea.
    [Notes:] Additional notes here
    
    Important guidelines:
    - Do not give only placeholders or labels. Write full, rich content for each bullet point.
    - Apply the same for the Questionnaire slide too, don't just keep placeholders, generate 10 questions.
    - Avoid general outlines — generate fully written content as if it's going directly into a slide.
    - Each bullet point should be 2-3 lines (30-50 words)
    - Provide complete thoughts with supporting details
    - Maintain parallel structure across points
    - Avoid single-sentence bullet points
    - Focus on substance over brevity
    
    Include these sections (adjust based on requested slide count):
    - Title slide
    - Introduction/Overview
    - Key Concepts
    - Detailed Analysis
    - Case Studies/Examples
    - Applications
    - Challenges
    - Future Trends
    - Conclusion

    After generating the content for the above presentation sections, add a clear separator:
    --- QUESTIONNAIRE ---
    
    Then generate 10 multiple-choice questions testing the user's understanding of the presentation:
    Format each question like:
    [Question: Your question text here?]
    - [A] Option A
    - [B] Option B
    - [C] Option C
    - [D] Option D
    [Correct: C]
    
    Ensure:
    - Questions cover all key areas: introduction, concepts, analysis, case studies, applications, challenges, and future trends
    - Each question has exactly 4 options
    - Clearly indicate the correct answer
    - Questions test comprehension of the presentation content

    Finally, add a separator:
    --- ANSWER KEY ---
    
    Then provide an answer key listing all correct answers in the format:
    [Answer Key]
    1. C
    2. A
    3. B
    ... [continue for all 10 questions]
    
    Please ensure the presentation content is presented first, followed by the questionnaire, then the answer key.
    Begin the content generation now.
    """
    
    response = model.generate_content(prompt)
    return response.text

def parse_slide_content(slide_text):
    slides = []
    current_slide = {}
    questionnaire = []
    answer_key = []

    # Split into sections
    parts = slide_text.split('--- QUESTIONNAIRE ---')
    main_content = parts[0].strip()
    remaining = parts[1].strip() if len(parts) > 1 else ""
    
    # Split remaining into questionnaire and answer key
    qa_parts = remaining.split('--- ANSWER KEY ---')
    questionnaire_content = qa_parts[0].strip() if qa_parts else ""
    answer_key_content = qa_parts[1].strip() if len(qa_parts) > 1 else ""
    
    # Parse main slides
    for line in slide_text.split('\n'):
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('[Title:]'):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'title': line.replace('[Title:]', '').strip(),
                'content': [],
                'notes': '',
                'layout': 'title-content'
            }
        elif line.startswith('[Content:]'):
            content = line.replace('[Content:]', '').strip()
            if content:
                current_slide['content'].append(content)
        elif line.startswith('[Notes:]'):
            current_slide['notes'] = line.replace('[Notes:]', '').strip()
        elif line.startswith('[Layout:]'):
            layout = line.replace('[Layout:]', '').strip().lower()
            valid_layouts = ['title-only', 'title-content', 'two-column', 'section-header']
            current_slide['layout'] = layout if layout in valid_layouts else 'title-content'
        elif current_slide.get('content') is not None and line.startswith('-'):
            # Simplify bullet points - remove any explanations after colons
            point = line[1:].strip()
            if ':' in point:
                point = point.split(':')[0].strip()
            current_slide['content'].append(point)
    
    if current_slide:
        slides.append(current_slide)

    # Parse questionnaire if exists
    if questionnaire_content:
        current_question = {}
        for line in questionnaire_content.split('\n'):
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('[Question:'):
                if current_question:
                    questionnaire.append(current_question)
                question_text = line.replace('[Question:', '').replace(']', '').strip()
                current_question = {
                    'text': question_text,
                    'options': [],
                    'correct': ''
                }
            elif line.startswith('- [A]'):
                current_question['options'].append(line.replace('- [A]', '').strip())
            elif line.startswith('- [B]'):
                current_question['options'].append(line.replace('- [B]', '').strip())
            elif line.startswith('- [C]'):
                current_question['options'].append(line.replace('- [C]', '').strip())
            elif line.startswith('- [D]'):
                current_question['options'].append(line.replace('- [D]', '').strip())
            elif line.startswith('[Correct:'):
                current_question['correct'] = line.replace('[Correct:', '').replace(']', '').strip()
        
        if current_question:
            questionnaire.append(current_question)

    # Parse answer key if exists
    if answer_key_content:
        # Extract answer key lines
        for line in answer_key_content.split('\n'):
            line = line.strip()
            if line.startswith('[') or not line:
                continue
            if re.match(r'^\d+\.\s*[A-D]$', line):
                answer_key.append(line.split('.')[1].strip())
    
    return slides, questionnaire, answer_key

def create_question_slide(prs, question, question_num, theme):
    """Create a slide for a single question"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    
    # Set title
    title = slide.shapes.title
    title.text = f"Question {question_num}"
    
    # Only apply custom formatting if not using a template
    if "template_path" not in theme:
        title.text_frame.paragraphs[0].font.color.rgb = theme["accent"]
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        if "title_font" in theme:
            title.text_frame.paragraphs[0].font.name = theme["title_font"]
    
    # Create content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    # Add question text
    p = tf.add_paragraph()
    p.text = question['text']
    p.level = 0
    p.font.bold = True
    p.space_after = Pt(24)
    
    # Only apply custom formatting if not using a template
    if "template_path" not in theme:
        p.font.color.rgb = theme["text_color"]
        p.font.size = Pt(24)
        if "text_font" in theme:
            p.font.name = theme["text_font"]
    
    # Add options
    for j, opt in enumerate(question['options']):
        p = tf.add_paragraph()
        p.text = f"{chr(65+j)}. {opt}"
        p.level = 0
        p.space_after = Pt(12)
        
        # Only apply custom formatting if not using a template
        if "template_path" not in theme:
            p.font.color.rgb = theme["text_color"]
            p.font.size = Pt(20)
            if "text_font" in theme:
                p.font.name = theme["text_font"]
    
    # Add correct answer to speaker notes
    notes_slide = slide.notes_slide
    notes_text = f"Correct Answer: {question['correct']}\n"
    notes_slide.notes_text_frame.text = notes_text
    
    return slide

def create_answer_key_slide(prs, questions, theme):
    """Create answer key slide with all correct answers"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
    
    # Set title
    title = slide.shapes.title
    title.text = "Answer Key"
    
    # Only apply custom formatting if not using a template
    if "template_path" not in theme:
        title.text_frame.paragraphs[0].font.color.rgb = theme["accent"]
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        if "title_font" in theme:
            title.text_frame.paragraphs[0].font.name = theme["title_font"]
    
    # Create content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    # Add answer key
    for i, q in enumerate(questions):
        p = tf.add_paragraph()
        p.text = f"Question {i+1}: {q['correct']}"
        p.level = 0
        p.space_after = Pt(12)
        
        # Only apply custom formatting if not using a template
        if "template_path" not in theme:
            p.font.color.rgb = theme["text_color"]
            p.font.size = Pt(24)
            if "text_font" in theme:
                p.font.name = theme["text_font"]
    
    return slide
    

def create_detailed_pptx(slides_data, questions, theme, branding_options=None):
    """Create PowerPoint using the uploaded template"""
    # Use the template if one was uploaded
    if "template_path" in theme and os.path.exists(theme["template_path"]):
        prs = Presentation(theme["template_path"])
    else:
        prs = Presentation()
    
    # Set widescreen layout (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Layout mapping
    layout_indices = {
        'title-only': 0,
        'title-content': 1,
        'section-header': 2,
        'two-column': 3
    }
    
    available_layouts = {}
    for name, idx in layout_indices.items():
        if idx < len(prs.slide_layouts):
            available_layouts[name] = idx
    
    default_layout_idx = available_layouts.get('title-content', 0)
    
    # Create main slides
    for slide_info in slides_data:
        layout = slide_info.get('layout', 'title-content').lower()
        layout_idx = available_layouts.get(layout, default_layout_idx)
        
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        except IndexError:
            slide = prs.slides.add_slide(prs.slide_layouts[default_layout_idx])
        
        # Set title
        title = slide.shapes.title
        title.text = slide_info['title']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Only apply custom formatting if not using a template
        if "template_path" not in theme:
            # Apply background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = theme["background"]
            
            # Format title
            title.text_frame.paragraphs[0].font.color.rgb = theme["title_color"]
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
            if "title_font" in theme:
                title.text_frame.paragraphs[0].font.name = theme["title_font"]
        
        # Add logo if provided
        if branding_options and branding_options.get('logo_path'):
            add_logo_to_slide(slide, branding_options['logo_path'], 
                            branding_options.get('logo_position', 'top-right'))
        
        # Set content based on layout
        if layout_idx == 3:  # Two column layout
            content = slide_info.get('content', [])
            mid_point = len(content) // 2
            left_content = content[:mid_point]
            right_content = content[mid_point:]
            
            left_body = slide.placeholders[1]
            left_tf = left_body.text_frame
            left_tf.clear()
            
            right_body = slide.placeholders[2]
            right_tf = right_body.text_frame
            right_tf.clear()
            
            for content_part, tf in [(left_content, left_tf), (right_content, right_tf)]:
                for point in content_part:
                    p = tf.add_paragraph()
                    point_text = point.replace('- ', '').strip()
                    p.text = point_text
                    p.level = 0
                    p.alignment = PP_ALIGN.JUSTIFY 
                    if "template_path" not in theme:  # Only apply custom formatting if no template
                        p.font.color.rgb = theme["text_color"]
                        p.font.size = Pt(18)
                        if "text_font" in theme:
                            p.font.name = theme["text_font"]
        
        elif layout_idx != 0:  # Not title-only
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            
            for point in slide_info.get('content', []):
                p = tf.add_paragraph()
                point_text = point.replace('- ', '').strip()
                p.text = point_text
                p.level = 0
                p.alignment = PP_ALIGN.JUSTIFY 
                if "template_path" not in theme:  # Only apply custom formatting if no template
                    p.font.color.rgb = theme["text_color"]
                    p.font.size = Pt(18)
                    if "text_font" in theme:
                        p.font.name = theme["text_font"]
        
        # Add notes if available
        if slide_info.get('notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_info['notes']
    
    # Add question slides (one per question)
    if questions:
        # Add section header
        section_slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section header layout
        section_title = section_slide.shapes.title
        section_title.text = "Knowledge Check"
        
        # Add each question on a separate slide
        for i, question in enumerate(questions, 1):
            create_question_slide(prs, question, i, theme)
        
        # Add answer key slide
        create_answer_key_slide(prs, questions, theme)
    
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    # Clean up temporary template file
    if "template_path" in theme and os.path.exists(theme["template_path"]):
        os.unlink(theme["template_path"])
    
    return pptx_io

def main():
    st.set_page_config(page_title="Advanced PPTX Generator", layout="wide")
    st.title("Advanced PowerPoint Generator")
    st.markdown("Create professional presentations with AI")

    # Initialize session state
    if 'custom_themes' not in st.session_state:
        st.session_state.custom_themes = {}
    if 'standardized_file' not in st.session_state:
        st.session_state.standardized_file = None
    if 'source_filename' not in st.session_state:
        st.session_state.source_filename = None

    # Combine themes
    ALL_THEMES = {**DEFAULT_THEMES, **st.session_state.custom_themes}

    # Tabs
    tab1, tab2 = st.tabs(["Generate New Presentation", "Standardize Existing Document"])

    with tab1:
        col1, col2 = st.columns([3, 1])

        with col1:
            topic = st.text_input("Presentation Topic:", placeholder="Enter your topic (e.g., 'AI in Healthcare')")

            with st.expander("Advanced Options"):
                slide_count = st.slider("Number of Slides:", 5, 20, 10)

                theme_option = st.radio("Theme Selection Method:", 
                                        ["Predefined Theme", "Custom Theme", "Example-Based Theme"])

                theme = DEFAULT_THEMES["Professional Blue"]
                uploaded_file = None

                if theme_option == "Predefined Theme":
                    theme_name = st.selectbox("Select Theme:", list(DEFAULT_THEMES.keys()))
                    theme = DEFAULT_THEMES[theme_name]

                elif theme_option == "Custom Theme":
                    theme_name = st.selectbox("Select Custom Theme:", ["Create New..."] + list(st.session_state.custom_themes.keys()))

                    if theme_name == "Create New...":
                        with st.form("custom_theme_form"):
                            new_theme_name = st.text_input("Theme Name")
                            bg_color = st.color_picker("Background Color", "#0C2340")
                            title_color = st.color_picker("Title Color", "#FFFFFF")
                            text_color = st.color_picker("Text Color", "#C8C8C8")
                            accent_color = st.color_picker("Accent Color", "#0070C0")
                            title_font = st.text_input("Title Font", "Calibri")
                            text_font = st.text_input("Text Font", "Calibri")

                            if st.form_submit_button("Save Custom Theme"):
                                if new_theme_name:
                                    st.session_state.custom_themes[new_theme_name] = {
                                        "background": hex_to_rgb(bg_color),
                                        "title_color": hex_to_rgb(title_color),
                                        "text_color": hex_to_rgb(text_color),
                                        "accent": hex_to_rgb(accent_color),
                                        "title_font": title_font,
                                        "text_font": text_font
                                    }
                                    st.success(f"Theme '{new_theme_name}' saved successfully!")
                                else:
                                    st.warning("Please enter a theme name")
                    elif theme_name in st.session_state.custom_themes:
                        theme = st.session_state.custom_themes[theme_name]

                else:  # Example-Based Theme
                    uploaded_file = st.file_uploader("Upload PowerPoint Template", type=["pptx"])
                    if uploaded_file:
                        with st.spinner("Extracting theme from template..."):
                            theme = extract_theme_from_pptx(uploaded_file)
                            st.success("Theme extracted from template!")
                    else:
                        st.info("Please upload a PowerPoint file to extract its theme")

        with col2:
            if st.button("Generate Presentation", type="primary", key="generate_btn"):
                if not topic:
                    st.warning("Please enter a topic first!")
                elif theme_option == "Example-Based Theme" and not uploaded_file:
                    st.warning("Please upload a PowerPoint template file first")
                else:
                    with st.spinner(f"Creating {slide_count}-slide presentation about '{topic}'..."):
                        try:
                            slide_text = generate_slide_content(topic, slide_count)
                            slides_data, questionnaire, answer_key = parse_slide_content(slide_text)

                            # Slide Overview
                            with st.expander("Slide Overview (Detailed)"):
                                for i, slide in enumerate(slides_data, 1):
                                    st.subheader(f"Slide {i}: {slide['title']}")
                                    st.markdown("**Content:**")
                                    for point in slide.get('content', []):
                                        st.markdown(f"- {point}")
                                    if slide.get('notes'):
                                        st.markdown(f"**Notes:** {slide['notes']}")
                                    st.markdown("---")
                            
                            # Create PPTX
                            if theme_option == "Example-Based Theme" and uploaded_file:
                                template_prs = Presentation(uploaded_file)
                                prs = Presentation()

                                # Remove default slide
                                blank_slide_layout = prs.slide_layouts[0]
                                prs.slides.add_slide(blank_slide_layout)

                                for slide in slides_data:
                                    slide_layout = template_prs.slide_layouts[1]
                                    new_slide = prs.slides.add_slide(slide_layout)

                                    title = new_slide.shapes.title
                                    title.text = slide['title']

                                    content = slide.get('content', [])
                                    if len(new_slide.placeholders) > 1:
                                        body = new_slide.placeholders[1]
                                        tf = body.text_frame
                                        tf.text = ""
                                        for point in content:
                                            p = tf.add_paragraph()
                                            p.text = point
                                            p.level = 0

                                pptx_file = io.BytesIO()
                                prs.save(pptx_file)
                                pptx_file.seek(0)

                            else:
                                pptx_file = create_detailed_pptx(slides_data, questionnaire, theme)

                            st.success("Presentation generated successfully!")

                            st.download_button(
                                label="Download PowerPoint",
                                data=pptx_file,
                                file_name=f"{topic.replace(' ', '_')}_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )

                        except Exception as e:
                            st.error(f"An error occurred: {str(e)}")                      

    with tab2:
        
        # ========== GENERALIZED STANDARDIZER FUNCTION ==========
        def standardize_ppt(source_path, template_path, output_dir=".", layout_index=1):
            prs_src = Presentation(source_path)
            prs_template = Presentation(template_path)
        
            # Remove all existing slides from template
            while len(prs_template.slides) > 0:
                r_id = prs_template.slides._sldIdLst[0].rId
                prs_template.part.drop_rel(r_id)
                del prs_template.slides._sldIdLst[0]
        
            layout = prs_template.slide_layouts[layout_index]
        
            for slide in prs_src.slides:
                new_slide = prs_template.slides.add_slide(layout)
        
                # --- Title ---
                title_text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        title_text = shape.text.strip()
                        break
        
                if new_slide.shapes.title:
                    title_shape = new_slide.shapes.title
                    title_shape.text_frame.clear()
                    title_shape.text = title_text
                    title_shape.text_frame.word_wrap = True
                    title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    for p in title_shape.text_frame.paragraphs:
                        p.font.size = Pt(24)
        
                # --- Body Placeholder or Fallback ---
                body_shape = None
        
                # 1. Try to find BODY placeholder
                for shape in new_slide.placeholders:
                    if shape.placeholder_format.type == 2:
                        body_shape = shape
                        break
        
                # 2. Else try other non-title placeholder with text_frame
                if not body_shape:
                    for shape in new_slide.placeholders:
                        if shape != new_slide.shapes.title and shape.has_text_frame:
                            body_shape = shape
                            break
        
                # 3. Else fallback: manually add a textbox
                if not body_shape:
                    body_shape = new_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        
                tf = body_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
                for src_shape in slide.shapes:
                    if src_shape.has_text_frame and src_shape.text.strip():
                        for para in src_shape.text_frame.paragraphs:
                            new_p = tf.add_paragraph()
                            new_p.alignment = para.alignment or PP_PARAGRAPH_ALIGNMENT.LEFT
                            for run in para.runs:
                                new_run = new_p.add_run()
                                new_run.text = run.text
                                new_run.font.size = run.font.size or Pt(16)
                                new_run.font.bold = run.font.bold
                                new_run.font.italic = run.font.italic
        
                                color = run.font.color
                                if color and hasattr(color, "rgb") and color.rgb:
                                    new_run.font.color.rgb = color.rgb
                
                # === PART 1: Copy standard header/footer placeholders ===
                for source in [prs_template.slide_master, prs_template.slide_layouts[layout_index]]:
                    for ph in source.placeholders:
                        fmt_type = ph.placeholder_format.type
                        if fmt_type in [11, 12, 13]:  # header, footer, slide number
                            try:
                                tf = ph.text_frame
                                if tf and tf.text.strip():
                                    box = new_slide.shapes.add_textbox(ph.left, ph.top, ph.width, ph.height)
                                    box_tf = box.text_frame
                                    box_tf.text = tf.text
                                    font = tf.paragraphs[0].font
                                    box_tf.paragraphs[0].font.size = font.size
                                    box_tf.paragraphs[0].font.name = font.name
                                    if font.color and hasattr(font.color, "rgb") and font.color.rgb:
                                        box_tf.paragraphs[0].font.color.rgb = font.color.rgb
                            except Exception as e:
                                print(f"Placeholder footer/header copy failed: {e}")
                
                # === PART 2: Copy manually placed shapes near bottom of slide (non-placeholders) ===
                for shape in prs_template.slide_layouts[layout_index].shapes:
                    if shape.has_text_frame and shape.text.strip() and shape.top > Inches(5):  # near bottom
                        try:
                            # Duplicate the shape at same position/dimension
                            new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                            new_tf = new_shape.text_frame
                            new_tf.clear()
                            new_tf.word_wrap = True
                
                            for para in shape.text_frame.paragraphs:
                                new_para = new_tf.add_paragraph()
                                new_para.alignment = para.alignment
                                for run in para.runs:
                                    new_run = new_para.add_run()
                                    new_run.text = run.text
                                    new_run.font.size = run.font.size
                                    new_run.font.name = run.font.name
                                    new_run.font.bold = run.font.bold
                                    new_run.font.italic = run.font.italic
                
                                    color = run.font.color
                                    if color and hasattr(color, "rgb") and color.rgb:
                                        new_run.font.color.rgb = color.rgb
                        except Exception as e:
                            print(f"⚠️ Failed to copy manual footer shape: {e}")
                            
            # Save result
            output_path = os.path.join(output_dir, f"standardized_{uuid.uuid4().hex}.pptx")
            prs_template.save(output_path)
            return output_path
        
        # ========== STREAMLIT UI ==========
        uploaded_source = st.file_uploader("📤 Upload Source PPTX", type=["pptx"])
        uploaded_template = st.file_uploader("🎨 Upload Template PPTX", type=["pptx"])
        
        if uploaded_source and uploaded_template:
            if st.button("🔁 Standardize Now"):
                # Save temporary files
                src_path = f"/tmp/source_{uuid.uuid4().hex}.pptx"
                tmpl_path = f"/tmp/template_{uuid.uuid4().hex}.pptx"
        
                with open(src_path, "wb") as f: f.write(uploaded_source.read())
                with open(tmpl_path, "wb") as f: f.write(uploaded_template.read())
        
                # Process standardization
                st.info("⚙️ Processing. Please wait...")
                output_path = standardize_ppt(src_path, tmpl_path, output_dir="/tmp")
        
                with open(output_path, "rb") as f:
                    st.success("✅ Standardization complete!")
                    st.download_button(
                        label="⬇ Download Standardized PPTX",
                        data=f,
                        file_name="standardized_output.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

if __name__ == "__main__":
    main()