import os
from pdfminer.high_level import extract_text
import docx2txt
from pptx import Presentation
import re

# Create output folders
TEXT_OUT = "paragraphs_txt"
LABEL_OUT = "paragraphs_labels"
os.makedirs(TEXT_OUT, exist_ok=True)
os.makedirs(LABEL_OUT, exist_ok=True)

def clean_paragraphs(text):
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]
    return paragraphs

def extract_from_pdf(path):
    return extract_text(path)

def extract_from_docx(path):
    return docx2txt.process(path)

def extract_from_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n\n".join(text)

def save_paragraphs(doc_name, paragraphs):
    for i, para in enumerate(paragraphs):
        txt_filename = f"{doc_name}_para{i}.txt"
        label_filename = f"{doc_name}_para{i}.label"

        with open(os.path.join(TEXT_OUT, txt_filename), "w", encoding="utf-8") as f:
            f.write(para)
        with open(os.path.join(LABEL_OUT, label_filename), "w", encoding="utf-8") as f:
            f.write("UNKNOWN")  # placeholder label

def process_file(file_path, file_name):
    try:
        if file_name.endswith(".pdf"):
            text = extract_from_pdf(file_path)
        elif file_name.endswith(".docx"):
            text = extract_from_docx(file_path)
        elif file_name.endswith(".pptx"):
            text = extract_from_pptx(file_path)
        else:
            print(f"Unsupported file format: {file_name}")
            return
        paragraphs = clean_paragraphs(text)
        doc_name = os.path.splitext(file_name)[0]
        save_paragraphs(doc_name, paragraphs)
        print(f"Processed: {file_name} ({len(paragraphs)} paragraphs)")
    except Exception as e:
        print(f"Error processing {file_name}: {e}")

# Main folder with your dataset files
DATASET_FOLDER = "Datasets"  # Replace with your folder path

for file in os.listdir(DATASET_FOLDER):
    process_file(os.path.join(DATASET_FOLDER, file), file)
